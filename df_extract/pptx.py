import json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from df_extract.base import BaseExtract
from df_extract.extract_util import cleanup
from df_extract.utils import iter_to_aiter


class ExtractPPTx(BaseExtract):
    """
    Class for PPTX or PPT Extraction
    """

    @staticmethod
    async def __extract_table(table):
        table_data = []
        async for _row in iter_to_aiter(table.rows):
            row_cols = []
            async for _cell in iter_to_aiter(_row.cells):
                row_cols.append(_cell.text)
            table_data.append(row_cols)
        return table_data

    @staticmethod
    async def __extract_text(text_frame):
        _paragraph = ""
        async for paragraph in iter_to_aiter(text_frame.paragraphs):
            async for run in iter_to_aiter(paragraph.runs):
                _paragraph += run.text
        return _paragraph

    async def _extract_shape(self, shape) -> str:
        text = ""
        _shape_type = shape.shape_type
        if _shape_type == MSO_SHAPE_TYPE.GROUP:
            _g_text = ""
            async for _shape in iter_to_aiter(shape.shapes):
                _g_text += await self._extract_shape(_shape)
            if _g_text:
                text += _g_text + "\n\n"
        elif _shape_type == MSO_SHAPE_TYPE.PICTURE:
            if self._image_extract:
                try:
                    image_date = await self.extract_image(shape.image.blob)
                    if image_date:
                        text += json.dumps(image_date) + "\n\n"
                except AttributeError as ex:
                    print(ex)
        elif _shape_type == MSO_SHAPE_TYPE.CHART:
            pass
        elif _shape_type == MSO_SHAPE_TYPE.TABLE:
            table_data = await self.__extract_table(shape.table)
            if table_data:
                text += json.dumps(table_data) + "\n\n"
        else:
            if shape.has_text_frame:
                # print("Shap Type => ", _shape.shape_type)
                # paragraph = ""
                # async for _paragraph in iter_to_aiter(_shape.text_frame.paragraphs):
                #     paragraph += _paragraph.text
                # print(paragraph)
                # paragraph = await self.extract_text(_shape.text_frame)
                # text += paragraph + "\n\n"
                _text = shape.text
                if _text:
                    text += _text + "\n\n"
        return text

    async def _extract_shapes(self, shapes):
        shapes_data = ""
        async for _shape in iter_to_aiter(shapes):
            _shape_type = _shape.shape_type
            if _shape.has_chart:
                pass
            if _shape.has_text_frame:
                pass
            if _shape.has_table:
                pass
            shapes_data += await self._extract_shape(_shape)
        cleaned_up_data = await cleanup(shapes_data)
        return cleaned_up_data

    async def extract_as_text(self, presentation):
        """
        Method to extract pptx or ppt content as text

        :param presentation: Valid `pptx.Presentation` object
        """
        await self.remove_existing_output()
        text = ""
        async for _slide in iter_to_aiter(presentation.slides):
            text += await self._extract_shapes(_slide.shapes)
            text += "\n\n\n\n"

        await self._write_text_output(text=text)

    async def extract_as_json(self, presentation):
        """
        Method to extract pptx or ppt content as json

        :param presentation: Valid `pptx.Presentation` object
        """
        await self.remove_existing_json_output()
        data = []
        slide_no = 1
        async for _slide in iter_to_aiter(presentation.slides):
            page_data = await self._extract_shapes(_slide.shapes)
            page_content = {
                'number': slide_no,
                'content': page_data,
                'name': self.name
            }
            data.append(page_content)
            slide_no += 1

        await self._write_json_output(data=data)

    async def extract(self, as_json: bool = False):
        """
        Method to extract csv content

        :param as_json: Extracted content will be stored as json. Default `False`
        """
        print(f'Extracting => {self.file_path}')
        presentation = Presentation(pptx=self.file_path)
        if not as_json:
            await self.extract_as_text(presentation)
        else:
            await self.extract_as_json(presentation)
        print(f'Extracted => {self.file_path}')
